"""
LocalNotebook v2
- チャットセッション永続化（複数セッション対応）
- 画像取り込み（Ollama vision で OCR+説明 → インデックス）
- ノートブック名リネーム
"""

import os, json, hashlib, re, time, base64
from pathlib import Path
from typing import Optional, List
from datetime import datetime
import asyncio

from fastapi import FastAPI, UploadFile, File, HTTPException, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, JSONResponse
from pydantic import BaseModel
import httpx

# -------- optional deps --------
try:
    import chromadb
    from chromadb.config import Settings
    CHROMA_OK = True
except ImportError:
    CHROMA_OK = False

try:
    from pypdf import PdfReader
    PDF_OK = True
except ImportError:
    PDF_OK = False

try:
    from docx import Document as DocxDocument
    DOCX_OK = True
except ImportError:
    DOCX_OK = False

try:
    import openpyxl
    XLSX_OK = True
except ImportError:
    XLSX_OK = False

try:
    from pptx import Presentation as PptxPresentation
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

# ================================
# Config
# ================================
OLLAMA_BASE  = os.getenv("OLLAMA_BASE",  "http://localhost:11434")
EMBED_MODEL  = os.getenv("EMBED_MODEL",  "nomic-embed-text")
CHAT_MODEL   = os.getenv("CHAT_MODEL",   "gemma3:4b")
VISION_MODEL = os.getenv("VISION_MODEL", "")   # 空なら自動検出
DATA_DIR     = Path(os.getenv("DATA_DIR", "./notebook_data"))
DATA_DIR.mkdir(parents=True, exist_ok=True)

NOTEBOOKS_FILE = DATA_DIR / "notebooks.json"
CHUNK_SIZE   = 400
CHUNK_OVERLAP = 80
IMAGE_EXTS   = {".png", ".jpg", ".jpeg", ".webp", ".gif"}
TABLE_EXTS   = {".xlsx", ".xls", ".xlsm", ".csv", ".tsv"}
PPTX_EXTS    = {".pptx", ".ppt"}

# ================================
# Storage helpers
# ================================
def load_notebooks() -> dict:
    if NOTEBOOKS_FILE.exists():
        return json.loads(NOTEBOOKS_FILE.read_text(encoding="utf-8"))
    return {}

def save_notebooks(data: dict):
    NOTEBOOKS_FILE.write_text(
        json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")

def get_nb(nb_id: str) -> dict:
    nbs = load_notebooks()
    if nb_id not in nbs:
        raise HTTPException(404, f"Notebook {nb_id} not found")
    return nbs[nb_id]

# ================================
# ChromaDB
# ================================
chroma_client = None
def get_chroma():
    global chroma_client
    if chroma_client is None and CHROMA_OK:
        chroma_client = chromadb.PersistentClient(
            path=str(DATA_DIR / "chroma"),
            settings=Settings(anonymized_telemetry=False))
    return chroma_client

def get_collection(nb_id: str):
    client = get_chroma()
    if client is None:
        return None
    safe = "nb_" + nb_id.replace("-", "_")
    return client.get_or_create_collection(safe, metadata={"nb_id": nb_id})

# ================================
# Ollama helpers
# ================================
async def ollama_embed(text: str) -> Optional[list]:
    async with httpx.AsyncClient(timeout=60) as c:
        try:
            r = await c.post(f"{OLLAMA_BASE}/api/embed",
                             json={"model": EMBED_MODEL, "input": text})
            data = r.json()
            embs = data.get("embeddings", [])
            return embs[0] if embs else None
        except Exception as e:
            print(f"[embed error] {e}")
            return None

async def ollama_chat_json(messages: list, model: str = None) -> str:
    """Non-streaming — Windows SSE 問題を回避"""
    m = model or CHAT_MODEL
    async with httpx.AsyncClient(timeout=180) as c:
        try:
            r = await c.post(f"{OLLAMA_BASE}/api/chat",
                             json={"model": m, "messages": messages, "stream": False})
            r.raise_for_status()
            return r.json().get("message", {}).get("content", "")
        except Exception as e:
            print(f"[chat error] {e}")
            raise HTTPException(500, f"Ollama chat error: {e}")

async def ollama_chat_stream(messages: list, model: str = None):
    """Streaming generator — summarize endpoint 用"""
    m = model or CHAT_MODEL
    async with httpx.AsyncClient(timeout=180) as c:
        async with c.stream("POST", f"{OLLAMA_BASE}/api/chat",
                            json={"model": m, "messages": messages, "stream": True}) as r:
            async for line in r.aiter_lines():
                if not line.strip():
                    continue
                try:
                    obj = json.loads(line)
                    tok = obj.get("message", {}).get("content", "")
                    if tok:
                        yield tok
                    if obj.get("done"):
                        break
                except Exception:
                    pass

async def ollama_vision(image_b64: str, media_type: str, model: str) -> str:
    """画像を vision モデルに渡してテキスト説明を得る"""
    prompt = ("この画像の内容を詳しく日本語で説明してください。"
              "文字が含まれる場合はすべて書き起こしてください。")
    async with httpx.AsyncClient(timeout=120) as c:
        try:
            r = await c.post(f"{OLLAMA_BASE}/api/chat",
                json={"model": model, "stream": False,
                      "messages": [{"role": "user", "content": prompt,
                                    "images": [image_b64]}]})
            r.raise_for_status()
            return r.json().get("message", {}).get("content", "")
        except Exception as e:
            print(f"[vision error] {e}")
            return ""

async def get_models() -> list:
    async with httpx.AsyncClient(timeout=10) as c:
        try:
            r = await c.get(f"{OLLAMA_BASE}/api/tags")
            return [m["name"] for m in r.json().get("models", [])]
        except Exception:
            return []

async def detect_vision_model() -> str:
    """vision 対応モデルを自動選択"""
    if VISION_MODEL:
        return VISION_MODEL
    models = await get_models()
    priority = ["llava", "gemma3", "minicpm-v", "moondream", "bakllava", "llama3.2-vision"]
    for p in priority:
        for m in models:
            if p in m.lower():
                return m
    return ""

# ================================
# Text extraction
# ================================
def extract_text(path: Path, mime: str = "") -> str:
    suffix = path.suffix.lower()
    if suffix in (".txt", ".md"):
        return path.read_text(encoding="utf-8", errors="replace")
    if suffix == ".pdf" and PDF_OK:
        reader = PdfReader(str(path))
        return "\n".join(p.extract_text() or "" for p in reader.pages)
    if suffix in (".docx", ".doc") and DOCX_OK:
        doc = DocxDocument(str(path))
        return "\n".join(p.text for p in doc.paragraphs)
    if suffix in TABLE_EXTS:
        return extract_table(path, suffix)
    if suffix in PPTX_EXTS:
        return extract_pptx(path)
    return path.read_text(encoding="utf-8", errors="replace")

def _rows_to_markdown(headers: list, rows: list) -> str:
    """行データをMarkdownテーブル形式に変換"""
    def fmt(v):
        return str(v).replace("|", "｜").replace("\n", " ") if v is not None else ""

    col_w = [max(len(h), max((len(fmt(r[i])) for r in rows if i < len(r)), default=0))
             for i, h in enumerate(headers)]
    def pad(s, w): return s.ljust(w)

    sep = "| " + " | ".join("-" * w for w in col_w) + " |"
    header_row = "| " + " | ".join(pad(h, col_w[i]) for i, h in enumerate(headers)) + " |"
    data_rows = []
    for row in rows:
        cells = [pad(fmt(row[i] if i < len(row) else ""), col_w[i]) for i in range(len(headers))]
        data_rows.append("| " + " | ".join(cells) + " |")
    return "\n".join([header_row, sep] + data_rows)


def _read_sheet_rows(ws) -> tuple:
    """シートから(headers, data_rows)を返す"""
    all_rows = []
    for row in ws.iter_rows():
        cells = [cell.value for cell in row]
        # 全空行はスキップ
        if any(v is not None and str(v).strip() for v in cells):
            all_rows.append(cells)
    if not all_rows:
        return [], []
    # 最初の非空行をヘッダーとして扱う
    headers = [str(v) if v is not None else f"列{i+1}" for i, v in enumerate(all_rows[0])]
    return headers, all_rows[1:]


def extract_table(path: Path, suffix: str) -> str:
    """Excel / CSV / TSV を LLM が理解しやすいMarkdown表形式に変換。
    ヘッダーを全チャンクに付与するため chunk_text_table() を利用。
    ここでは全体テキストを返す（チャンク分割は chunk_table_aware で行う）。
    """
    import csv as csv_mod

    sections = []  # (sheet_name, headers, data_rows) のリスト

    if suffix in (".csv", ".tsv"):
        delim = "\t" if suffix == ".tsv" else ","
        raw_rows = []
        for enc in ("utf-8", "utf-8-sig", "shift_jis", "cp932"):
            try:
                with open(path, encoding=enc, errors="strict", newline="") as f:
                    raw_rows = list(csv_mod.reader(f, delimiter=delim))
                break
            except (UnicodeDecodeError, Exception):
                continue
        if not raw_rows:
            return "[CSV読み取りエラー]"
        headers = [str(v) if v else f"列{i+1}" for i, v in enumerate(raw_rows[0])]
        sections.append((path.name, headers, raw_rows[1:]))

    elif XLSX_OK:
        wb = openpyxl.load_workbook(str(path), data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            headers, data_rows = _read_sheet_rows(ws)
            if headers:
                sections.append((sheet_name, headers, data_rows))
    else:
        return "[xlsx読み取りエラー: openpyxlが必要です]"

    if not sections:
        return "[データなし]"

    parts = []
    for sheet_name, headers, data_rows in sections:
        # シート概要
        summary = (f"=== {sheet_name} ===\n"
                   f"行数: {len(data_rows)}行 / 列数: {len(headers)}列\n"
                   f"列名: {', '.join(headers)}\n")
        parts.append(summary)
        # Markdownテーブル（全行）
        if data_rows:
            parts.append(_rows_to_markdown(headers, data_rows))
        parts.append("")

    return "\n".join(parts)


def chunk_table_aware(text: str, filename: str,
                      size: int = CHUNK_SIZE * 3,
                      overlap_rows: int = 2) -> List[str]:
    """
    表データ専用チャンク分割。
    ヘッダー行と概要を各チャンクの先頭に付与し、
    LLMがどのチャンクでも列の意味を把握できるようにする。
    """
    chunks = []
    # セクション（シート）ごとに分割
    sections = re.split(r"(=== .+ ===\n)", text)
    i = 0
    while i < len(sections):
        block = sections[i]
        if re.match(r"=== .+ ===\n", block) and i + 1 < len(sections):
            header_block = block + sections[i + 1]
            i += 2
        else:
            header_block = block
            i += 1
        if not header_block.strip():
            continue

        lines = header_block.split("\n")
        # Markdownテーブルのヘッダー行（| col1 | col2 | ...）を探す
        md_header_idx = None
        md_sep_idx = None
        for j, line in enumerate(lines):
            if line.startswith("|") and "---" not in line and md_header_idx is None:
                md_header_idx = j
            elif line.startswith("|") and "---" in line:
                md_sep_idx = j
                break

        if md_header_idx is None:
            # テーブル形式でない → そのまま通常チャンク
            chunks.extend(chunk_text(header_block))
            continue

        prefix_lines = lines[:md_sep_idx + 1] if md_sep_idx else lines[:md_header_idx + 1]
        prefix = "\n".join(prefix_lines) + "\n"
        data_lines = lines[md_sep_idx + 1:] if md_sep_idx else lines[md_header_idx + 1:]
        data_lines = [l for l in data_lines if l.strip()]

        # データ行を size に収まるよう分割（ヘッダーを毎回付ける）
        chunk_rows = []
        for row_line in data_lines:
            chunk_rows.append(row_line)
            candidate = prefix + "\n".join(chunk_rows)
            if len(candidate) >= size:
                chunks.append(candidate)
                # overlap: 最後の overlap_rows 行を次のチャンクに引き継ぐ
                chunk_rows = chunk_rows[-overlap_rows:]
        if chunk_rows:
            chunks.append(prefix + "\n".join(chunk_rows))

    return [c for c in chunks if c.strip()] or [text[:size]]

def extract_pptx(path: Path) -> str:
    """PowerPoint の各スライドからテキストを抽出"""
    if not PPTX_OK:
        return "[pptx読み取りエラー: python-pptxが必要です]"
    prs = PptxPresentation(str(path))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"=== スライド {i} ===")
        # スライドタイトル
        if slide.shapes.title and slide.shapes.title.text.strip():
            lines.append(f"タイトル: {slide.shapes.title.text.strip()}")
        # 全テキストボックス
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape == slide.shapes.title:
                continue
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if txt:
                    indent = "  " * para.level
                    lines.append(f"{indent}{txt}")
        # テーブル
        for shape in slide.shapes:
            if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                tbl = shape.table
                for row in tbl.rows:
                    cells = [c.text.strip() for c in row.cells]
                    lines.append("\t".join(cells))
        # ノート
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"[ノート] {notes}")
        lines.append("")
    return "\n".join(lines)

def chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    text = re.sub(r'\n{3,}', '\n\n', text).strip()
    chunks, start = [], 0
    while start < len(text):
        chunk = text[start:start + size]
        if chunk.strip():
            chunks.append(chunk)
        start += size - overlap
    return chunks

# ================================
# Index helpers
# ================================
async def index_document(nb_id: str, doc_id: str, text: str,
                          doc_type: str = "document"):
    col = get_collection(nb_id)
    if col is None:
        return
    try:
        existing = col.get(where={"doc_id": doc_id})
        if existing["ids"]:
            col.delete(ids=existing["ids"])
    except Exception:
        pass

    if doc_type == "table":
        chunks = chunk_table_aware(text, doc_id)
    else:
        chunks = chunk_text(text)
    ids, docs, metas, embs = [], [], [], []
    for i, chunk in enumerate(chunks):
        emb = await ollama_embed(chunk)
        if emb is None:
            continue
        ids.append(f"{doc_id}_chunk_{i}")
        docs.append(chunk)
        metas.append({"doc_id": doc_id, "chunk_idx": i})
        embs.append(emb)
    if ids:
        col.add(ids=ids, documents=docs, metadatas=metas, embeddings=embs)

# ================================
# Chat session helpers
# ================================
def sessions_file(nb_id: str) -> Path:
    p = DATA_DIR / nb_id
    p.mkdir(exist_ok=True)
    return p / "chat_sessions.json"

def load_sessions(nb_id: str) -> dict:
    f = sessions_file(nb_id)
    if f.exists():
        return json.loads(f.read_text(encoding="utf-8"))
    return {}

def save_sessions(nb_id: str, data: dict):
    sessions_file(nb_id).write_text(
        json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")

def new_session_id() -> str:
    return hashlib.md5(str(time.time()).encode()).hexdigest()[:10]

# ================================
# FastAPI app
# ================================
app = FastAPI(title="LocalNotebook v2")
app.add_middleware(CORSMiddleware, allow_origins=["*"],
                  allow_methods=["*"], allow_headers=["*"])

# ---- Notebooks CRUD ----
class NotebookCreate(BaseModel):
    name: str
    description: str = ""

@app.get("/api/notebooks")
def list_notebooks():
    return load_notebooks()

@app.post("/api/notebooks")
def create_notebook(body: NotebookCreate):
    nbs = load_notebooks()
    nb_id = hashlib.md5(f"{body.name}{time.time()}".encode()).hexdigest()[:12]
    nbs[nb_id] = {
        "id": nb_id, "name": body.name, "description": body.description,
        "created_at": datetime.now().isoformat(),
        "documents": {}, "notes": []
    }
    save_notebooks(nbs)
    return nbs[nb_id]

@app.delete("/api/notebooks/{nb_id}")
def delete_notebook(nb_id: str):
    nbs = load_notebooks()
    if nb_id not in nbs:
        raise HTTPException(404)
    del nbs[nb_id]
    save_notebooks(nbs)
    try:
        client = get_chroma()
        if client:
            client.delete_collection("nb_" + nb_id.replace("-", "_"))
    except Exception:
        pass
    return {"ok": True}

@app.patch("/api/notebooks/{nb_id}")
def update_notebook(nb_id: str, body: dict):
    nbs = load_notebooks()
    nb = nbs.get(nb_id)
    if not nb:
        raise HTTPException(404)
    for k in ("name", "description"):
        if k in body:
            nb[k] = body[k]
    save_notebooks(nbs)
    return nb

# ---- Documents ----
@app.post("/api/notebooks/{nb_id}/documents")
async def upload_document(nb_id: str, background_tasks: BackgroundTasks,
                          file: UploadFile = File(...)):
    nbs = load_notebooks()
    nb = nbs.get(nb_id)
    if not nb:
        raise HTTPException(404)

    content = await file.read()
    doc_id = hashlib.md5(content).hexdigest()[:12]
    suffix = Path(file.filename).suffix.lower()
    is_image = suffix in IMAGE_EXTS
    is_table = suffix in TABLE_EXTS
    is_pptx  = suffix in PPTX_EXTS

    doc_dir = DATA_DIR / nb_id
    doc_dir.mkdir(exist_ok=True)
    fpath = doc_dir / f"{doc_id}_{file.filename}"
    fpath.write_bytes(content)

    doc_meta = {
        "id": doc_id, "filename": file.filename,
        "size": len(content),
        "uploaded_at": datetime.now().isoformat(),
        "status": "indexing",
        "type": "image" if is_image else ("table" if is_table else ("pptx" if is_pptx else "document")),
        "word_count": 0,
        "text_preview": ""
    }
    nbs[nb_id]["documents"][doc_id] = doc_meta
    save_notebooks(nbs)

    async def do_index():
        text = ""
        try:
            if is_image:
                vision_m = await detect_vision_model()
                if vision_m:
                    b64 = base64.b64encode(content).decode()
                    desc = await ollama_vision(b64, file.content_type or "", vision_m)
                    text = f"【画像ファイル: {file.filename}】\n{desc}"
                else:
                    text = f"【画像ファイル: {file.filename}】（vision モデル未検出のためテキスト化スキップ）"
            elif is_table:
                text = extract_table(fpath, suffix)
                text = f"【表データ: {file.filename}】\n{text}"
            elif is_pptx:
                text = extract_pptx(fpath)
                text = f"【プレゼン: {file.filename}】\n{text}"
            else:
                text = extract_text(fpath, file.content_type or "")
        except Exception as e:
            text = f"[抽出エラー: {e}]"

        dtype = "table" if is_table else ("image" if is_image else ("pptx" if is_pptx else "document"))
        await index_document(nb_id, doc_id, text, doc_type=dtype)
        nbs2 = load_notebooks()
        if nb_id in nbs2 and doc_id in nbs2[nb_id]["documents"]:
            nbs2[nb_id]["documents"][doc_id]["status"] = "ready"
            nbs2[nb_id]["documents"][doc_id]["word_count"] = len(text.split())
            nbs2[nb_id]["documents"][doc_id]["text_preview"] = text[:300]
            save_notebooks(nbs2)

    background_tasks.add_task(do_index)
    return doc_meta

@app.delete("/api/notebooks/{nb_id}/documents/{doc_id}")
def delete_document(nb_id: str, doc_id: str):
    nbs = load_notebooks()
    nb = nbs.get(nb_id)
    if not nb or doc_id not in nb["documents"]:
        raise HTTPException(404)
    del nb["documents"][doc_id]
    save_notebooks(nbs)
    try:
        col = get_collection(nb_id)
        if col:
            existing = col.get(where={"doc_id": doc_id})
            if existing["ids"]:
                col.delete(ids=existing["ids"])
    except Exception:
        pass
    return {"ok": True}

# ---- Search ----
class SearchQuery(BaseModel):
    query: str
    n_results: int = 5

@app.post("/api/notebooks/{nb_id}/search")
async def search_documents(nb_id: str, body: SearchQuery):
    get_nb(nb_id)
    col = get_collection(nb_id)
    if not col:
        raise HTTPException(500, "ChromaDB not available")
    emb = await ollama_embed(body.query)
    if emb is None:
        raise HTTPException(500, "Embedding failed")
    results = col.query(
        query_embeddings=[emb],
        n_results=min(body.n_results, col.count() or 1),
        include=["documents", "metadatas", "distances"])
    nbs = load_notebooks()
    docs = nbs.get(nb_id, {}).get("documents", {})
    hits = []
    for doc, meta, dist in zip(results["documents"][0],
                                results["metadatas"][0],
                                results["distances"][0]):
        doc_info = docs.get(meta["doc_id"], {})
        hits.append({"chunk": doc, "doc_id": meta["doc_id"],
                     "filename": doc_info.get("filename", "unknown"),
                     "score": round(1 - dist, 3),
                     "chunk_idx": meta.get("chunk_idx", 0)})
    return {"results": hits}

# ---- Chat sessions ----
@app.get("/api/notebooks/{nb_id}/sessions")
def list_sessions(nb_id: str):
    get_nb(nb_id)
    sessions = load_sessions(nb_id)
    # メッセージ本体は省いてメタだけ返す
    return {sid: {k: v for k, v in s.items() if k != "messages"}
            for sid, s in sessions.items()}

@app.post("/api/notebooks/{nb_id}/sessions")
def create_session(nb_id: str, body: dict = {}):
    get_nb(nb_id)
    sessions = load_sessions(nb_id)
    sid = new_session_id()
    sessions[sid] = {
        "id": sid,
        "title": body.get("title", "新しいチャット"),
        "created_at": datetime.now().isoformat(),
        "updated_at": datetime.now().isoformat(),
        "messages": []
    }
    save_sessions(nb_id, sessions)
    return sessions[sid]

@app.get("/api/notebooks/{nb_id}/sessions/{sid}")
def get_session(nb_id: str, sid: str):
    get_nb(nb_id)
    sessions = load_sessions(nb_id)
    if sid not in sessions:
        raise HTTPException(404, "Session not found")
    return sessions[sid]

@app.delete("/api/notebooks/{nb_id}/sessions/{sid}")
def delete_session(nb_id: str, sid: str):
    get_nb(nb_id)
    sessions = load_sessions(nb_id)
    if sid not in sessions:
        raise HTTPException(404)
    del sessions[sid]
    save_sessions(nb_id, sessions)
    return {"ok": True}

@app.patch("/api/notebooks/{nb_id}/sessions/{sid}")
def rename_session(nb_id: str, sid: str, body: dict):
    get_nb(nb_id)
    sessions = load_sessions(nb_id)
    if sid not in sessions:
        raise HTTPException(404)
    if "title" in body:
        sessions[sid]["title"] = body["title"]
    save_sessions(nb_id, sessions)
    return {k: v for k, v in sessions[sid].items() if k != "messages"}

# ---- Q&A / Chat ----
class ChatMessage(BaseModel):
    role: str
    content: str

class ChatRequest(BaseModel):
    messages: List[ChatMessage]
    session_id: Optional[str] = None   # 指定があれば履歴に保存
    use_docs: bool = True
    model: Optional[str] = None

@app.post("/api/notebooks/{nb_id}/chat")
async def chat(nb_id: str, body: ChatRequest):
    get_nb(nb_id)
    user_msg = body.messages[-1].content

    # RAG
    context_chunks = []
    if body.use_docs:
        col = get_collection(nb_id)
        if col and col.count() > 0:
            emb = await ollama_embed(user_msg)
            if emb:
                results = col.query(
                    query_embeddings=[emb],
                    n_results=min(6, col.count()),
                    include=["documents", "metadatas"])
                nbs = load_notebooks()
                docs = nbs.get(nb_id, {}).get("documents", {})
                for doc, meta in zip(results["documents"][0], results["metadatas"][0]):
                    fn = docs.get(meta["doc_id"], {}).get("filename", "doc")
                    context_chunks.append(f"[{fn}]\n{doc}")

    system = "あなたは優秀なリサーチアシスタントです。日本語で回答してください。"
    if context_chunks:
        ctx = "\n\n---\n\n".join(context_chunks)
        system += f"\n\n以下のドキュメントのコンテキストを参考にして回答してください:\n\n{ctx}"

    messages = [{"role": "system", "content": system}]
    for m in body.messages[:-1][-10:]:
        messages.append({"role": m.role, "content": m.content})
    messages.append({"role": "user", "content": user_msg})

    reply = await ollama_chat_json(messages, body.model)

    # セッション保存
    if body.session_id:
        sessions = load_sessions(nb_id)
        if body.session_id in sessions:
            sess = sessions[body.session_id]
            # 既存履歴を body.messages で上書き（クライアントが正）
            sess["messages"] = [{"role": m.role, "content": m.content}
                                 for m in body.messages]
            sess["messages"].append({"role": "assistant", "content": reply})
            sess["updated_at"] = datetime.now().isoformat()
            # タイトルが「新しいチャット」のままなら最初のユーザー発言から自動生成
            if sess["title"] == "新しいチャット":
                sess["title"] = user_msg[:30] + ("…" if len(user_msg) > 30 else "")
            save_sessions(nb_id, sessions)

    return JSONResponse({"reply": reply})

# ---- Summarize ----
class SummarizeRequest(BaseModel):
    doc_id: Optional[str] = None
    style: str = "bullet"

@app.post("/api/notebooks/{nb_id}/summarize")
async def summarize(nb_id: str, body: SummarizeRequest):
    nbs = load_notebooks()
    nb = nbs.get(nb_id)
    if not nb:
        raise HTTPException(404)
    doc_dir = DATA_DIR / nb_id
    texts = {}
    if body.doc_id:
        doc = nb["documents"].get(body.doc_id)
        if not doc:
            raise HTTPException(404, "Document not found")
        for f in doc_dir.glob(f"{body.doc_id}_*"):
            texts[doc["filename"]] = extract_text(f)
    else:
        for doc_id, doc in nb["documents"].items():
            for f in doc_dir.glob(f"{doc_id}_*"):
                texts[doc["filename"]] = extract_text(f)[:3000]
    if not texts:
        raise HTTPException(400, "No text to summarize")
    style_prompt = {
        "bullet": "箇条書き形式で重要なポイントを10個以内にまとめてください。",
        "detailed": "詳細な要約を日本語で書いてください。各セクションに見出しをつけてください。",
        "executive": "エグゼクティブサマリーとして、1〜2段落で核心的な内容を簡潔にまとめてください。"
    }.get(body.style, "要約してください。")
    combined = "\n\n=====\n\n".join(f"【{fn}】\n{t}" for fn, t in texts.items())
    messages = [
        {"role": "system", "content": "あなたは優秀な文書要約アシスタントです。"},
        {"role": "user", "content": f"以下のドキュメントを{style_prompt}\n\n{combined[:6000]}"}
    ]
    async def generate():
        async for tok in ollama_chat_stream(messages):
            yield f"data: {json.dumps({'token': tok})}\n\n"
        yield "data: [DONE]\n\n"
    return StreamingResponse(generate(), media_type="text/event-stream")

# ---- Notes ----
class NoteCreate(BaseModel):
    content: str
    title: str = ""

@app.post("/api/notebooks/{nb_id}/notes")
def add_note(nb_id: str, body: NoteCreate):
    nbs = load_notebooks()
    nb = nbs.get(nb_id)
    if not nb:
        raise HTTPException(404)
    note = {
        "id": hashlib.md5(f"{body.content}{time.time()}".encode()).hexdigest()[:8],
        "title": body.title or body.content[:40],
        "content": body.content,
        "created_at": datetime.now().isoformat()
    }
    nb.setdefault("notes", []).append(note)
    save_notebooks(nbs)
    return note

@app.delete("/api/notebooks/{nb_id}/notes/{note_id}")
def delete_note(nb_id: str, note_id: str):
    nbs = load_notebooks()
    nb = nbs.get(nb_id)
    if not nb:
        raise HTTPException(404)
    nb["notes"] = [n for n in nb.get("notes", []) if n["id"] != note_id]
    save_notebooks(nbs)
    return {"ok": True}

# ---- Models / Status ----
@app.get("/api/models")
async def list_models():
    return {"models": await get_models()}

@app.get("/api/status")
async def status():
    models = await get_models()
    vision = await detect_vision_model()
    return {
        "ollama": len(models) > 0,
        "chroma": CHROMA_OK, "pdf": PDF_OK, "docx": DOCX_OK,
        "xlsx": XLSX_OK, "pptx": PPTX_OK,
        "models": models,
        "embed_model": EMBED_MODEL, "chat_model": CHAT_MODEL,
        "vision_model": vision
    }

if __name__ == "__main__":
    import uvicorn, sys
    if sys.platform == "win32":
        asyncio.set_event_loop_policy(asyncio.WindowsSelectorEventLoopPolicy())
    uvicorn.run(app, host="0.0.0.0", port=8765, reload=False)
